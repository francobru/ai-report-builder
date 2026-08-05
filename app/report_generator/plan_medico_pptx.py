"""PPTX generator for the Plan Medico monthly report.

Replicates the May 2026 Plan Medico template, minus the two sections the
user dropped ("Resolucion en llamado" and the No-Disponible breakdown).

Slides produced:
    1. Cover with the three headline figures
    2. Seccion 1 - Plan Medico Total
    3. Seccion 2 - Turnos Plan Medico (+ per-skill table)
    4. Seccion 3a - PM Consultas
    5. Seccion 3b - Cierres de llamados (top reasons)
    6. Seccion 3c - Productividad de Agentes
"""

from __future__ import annotations

import base64
import io
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE


# ======================================================================
# Palette (Plan Medico template)
# ======================================================================
BLUE = RGBColor(0x2F, 0x5F, 0xD0)
LIGHT_BLUE = RGBColor(0x4F, 0xC3, 0xF7)
GREEN = RGBColor(0x8B, 0xC3, 0x4A)
RED = RGBColor(0xF4, 0x43, 0x36)
PURPLE = RGBColor(0x7C, 0x4D, 0xFF)
CARD_BG = RGBColor(0xF4, 0xF7, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK = RGBColor(0x2B, 0x33, 0x45)
TEXT_GRAY = RGBColor(0x6B, 0x72, 0x80)
BORDER = RGBColor(0xE3, 0xE8, 0xF0)
ROW_ALT = RGBColor(0xF7, 0xF9, 0xFC)

FONT = "Calibri"
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ======================================================================
# Primitives
# ======================================================================

def _rect(slide, x, y, w, h, fill=None, line=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sh.shadow.inherit = False
    if fill:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    else:
        sh.fill.background()
    if line:
        sh.line.fill.solid()
        sh.line.fill.fore_color.rgb = line
        sh.line.width = Pt(0.75)
    else:
        sh.line.fill.background()
    return sh


def _text(slide, x, y, w, h, text, size=11, bold=False, color=TEXT_DARK,
          align=PP_ALIGN.LEFT, italic=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = FONT
    p.alignment = align
    return box


def _image(slide, src, x, y, w, h):
    """Insert an image preserving its aspect ratio."""
    if not src:
        return
    if isinstance(src, (str, Path)):
        p = Path(src)
        if not p.exists():
            return
        src = str(p)
    try:
        from PIL import Image as PILImage
        iw, ih = PILImage.open(src).size
        scale = min(w / (iw / 96.0), h / (ih / 96.0))
        dw, dh = (iw / 96.0) * scale, (ih / 96.0) * scale
        slide.shapes.add_picture(src, Inches(x + (w - dw) / 2),
                                 Inches(y + (h - dh) / 2), Inches(dw), Inches(dh))
    except Exception:
        slide.shapes.add_picture(src, Inches(x), Inches(y), Inches(w), Inches(h))


def _header(slide, title, subtitle):
    _rect(slide, 0, 0, 13.333, 0.06, fill=BLUE)
    _text(slide, 0.6, 0.32, 11.5, 0.45, title, size=22, bold=True, color=TEXT_DARK)
    if subtitle:
        _text(slide, 0.6, 0.78, 11.5, 0.3, subtitle, size=11, color=TEXT_GRAY)
    _rect(slide, 0.6, 1.16, 12.1, 0.012, fill=BORDER)


def _footer(slide, period, page, total):
    _text(slide, 0.6, 7.06, 6, 0.25, f"Reporte Plan Medico - {period}",
          size=8, color=TEXT_GRAY)
    _text(slide, 7.0, 7.06, 5.7, 0.25, f"Pagina {page} de {total}",
          size=8, color=TEXT_GRAY, align=PP_ALIGN.RIGHT)


def _kpi_card(slide, x, y, w, label, value, accent, h=1.05, value_size=26):
    """Card with a coloured top bar, big value and caption underneath."""
    _rect(slide, x, y, w, 0.11, fill=accent)
    _rect(slide, x, y + 0.11, w, h - 0.11, fill=CARD_BG, line=BORDER)
    _text(slide, x + 0.1, y + 0.30, w - 0.2, 0.45, value,
          size=value_size, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _text(slide, x + 0.1, y + h - 0.32, w - 0.2, 0.25, label,
          size=10, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)


def _table(slide, x, y, widths, headers, rows, header_fill=BLUE,
           row_h=0.32, header_h=0.34, font_size=10, header_size=10,
           bold_last=False, align_first_left=True):
    n_rows = len(rows) + 1
    total_w = sum(widths)
    shape = slide.shapes.add_table(n_rows, len(widths), Inches(x), Inches(y),
                                   Inches(total_w),
                                   Inches(header_h + row_h * len(rows)))
    tbl = shape.table
    for i, wd in enumerate(widths):
        tbl.columns[i].width = Inches(wd)
    tbl.rows[0].height = Inches(header_h)
    for i in range(len(rows)):
        tbl.rows[i + 1].height = Inches(row_h)

    for j, htxt in enumerate(headers):
        c = tbl.cell(0, j)
        c.text = str(htxt)
        c.fill.solid()
        c.fill.fore_color.rgb = header_fill
        for p in c.text_frame.paragraphs:
            p.font.size = Pt(header_size)
            p.font.bold = True
            p.font.name = FONT
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.LEFT if (j == 0 and align_first_left) else PP_ALIGN.RIGHT

    for i, row in enumerate(rows):
        last = bold_last and i == len(rows) - 1
        bg = ROW_ALT if (i % 2 == 0 or last) else WHITE
        for j, val in enumerate(row):
            c = tbl.cell(i + 1, j)
            c.text = str(val)
            c.fill.solid()
            c.fill.fore_color.rgb = bg
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.bold = bool(last)
                p.font.name = FONT
                p.font.color.rgb = BLUE if last else TEXT_DARK
                p.alignment = PP_ALIGN.LEFT if (j == 0 and align_first_left) else PP_ALIGN.RIGHT
    return shape


# ======================================================================
# Slides
# ======================================================================

def _slide_cover(prs, period, headline):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _rect(s, 0, 0, 13.333, 0.06, fill=BLUE)

    try:
        from app.report_generator.logo_data import LOGO_HA_APAISADO_B64
        stream = io.BytesIO(base64.b64decode(LOGO_HA_APAISADO_B64))
        w = 5.6
        s.shapes.add_picture(stream, Inches((13.333 - w) / 2), Inches(0.75),
                             Inches(w), Inches(w / 4.95))
    except Exception:
        pass

    _text(s, 0, 2.35, 13.333, 0.6, "Reporte Plan Medico", size=34, bold=True,
          color=TEXT_DARK, align=PP_ALIGN.CENTER)
    _text(s, 0, 2.98, 13.333, 0.4, "Contact Center", size=20,
          color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    _text(s, 0, 3.40, 13.333, 0.4, period, size=20, bold=True,
          color=BLUE, align=PP_ALIGN.CENTER)
    _text(s, 0, 3.84, 13.333, 0.3,
          "Visualizacion integral de gestion, eficiencia y productividad",
          size=12, color=TEXT_GRAY, align=PP_ALIGN.CENTER)

    cards = [
        ("PM Total\n(Recibidas)", headline.get("pm_total", "-"), BLUE),
        ("Turnos PM\n(Recibidas)", headline.get("turnos_pm", "-"), LIGHT_BLUE),
        ("PM Consultas\n(Recibidas)", headline.get("pm_consultas", "-"), PURPLE),
    ]
    cw, gap = 3.4, 0.5
    x0 = (13.333 - (cw * 3 + gap * 2)) / 2
    for i, (label, value, accent) in enumerate(cards):
        x = x0 + i * (cw + gap)
        _rect(s, x, 4.55, cw, 0.11, fill=accent)
        _rect(s, x, 4.66, cw, 1.35, fill=CARD_BG, line=BORDER)
        _text(s, x + 0.1, 4.90, cw - 0.2, 0.55, value, size=30, bold=True,
              color=accent, align=PP_ALIGN.CENTER)
        box = _text(s, x + 0.1, 5.48, cw - 0.2, 0.45, label.replace("\n", " "),
                    size=11, bold=True, color=TEXT_DARK, align=PP_ALIGN.CENTER)


def _slide_section(prs, period, page, total, title, subtitle, kpis, chart_path,
                   accent=BLUE, table=None):
    """A section slide: 5 KPI cards, an optional table, and the daily chart."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, title, subtitle)
    _footer(s, period, page, total)

    specs = [
        ("Llamadas Recibidas", kpis.get("recibidas", "-"), BLUE),
        ("Llamadas Atendidas", kpis.get("atendidas", "-"), LIGHT_BLUE),
        ("Nivel de Atencion", kpis.get("nivel_atencion", "-"), GREEN),
        ("No Atendidas", kpis.get("no_atendidas", "-"), RED),
        ("Conv. promedio", kpis.get("conversacion", "-"), PURPLE),
    ]
    cw, gap = 2.32, 0.15
    x0 = (13.333 - (cw * 5 + gap * 4)) / 2
    for i, (label, value, color) in enumerate(specs):
        _kpi_card(s, x0 + i * (cw + gap), 1.32, cw, label, value, color)

    y = 2.55
    if table:
        _table(s, table["x"], y, table["widths"], table["headers"], table["rows"],
               bold_last=table.get("bold_last", True))
        y += 0.34 + 0.32 * len(table["rows"]) + 0.25

    if chart_path:
        _image(s, chart_path, 0.5, y, 12.3, 6.85 - y)


def _slide_closures(prs, period, page, total, subtitle, chart_path):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "Seccion 3b - PM Consultas: Cierres de llamados", subtitle)
    _footer(s, period, page, total)
    _image(s, chart_path, 0.5, 1.35, 12.3, 5.5)


def _slide_agents(prs, period, page, total, summary, rows):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    _header(s, "Seccion 3c - PM Consultas: Productividad de Agentes",
            f"{summary['operadores']} operadores con atendidas y cierres "
            f"en la habilidad PM Consultas durante {period}")
    _footer(s, period, page, total)

    cards = [
        ("Tiempo total de logueo", summary["logueo"], BLUE),
        ("Conv. promedio", summary["aht"], PURPLE),
        (f"Tiempo No Disponible ({str(summary['pct_no_disponible']).replace('.', ',')}%)",
         summary["no_disponible"], RED),
    ]
    cw, gap = 3.6, 0.35
    x0 = (13.333 - (cw * 3 + gap * 2)) / 2
    for i, (label, value, color) in enumerate(cards):
        _kpi_card(s, x0 + i * (cw + gap), 1.30, cw, label, value, color,
                  h=1.0, value_size=22)

    headers = ["Operador", "Logueo", "No Disp.", "% No Disp.",
               "Atendidas", "Cierres", "Util.%", "Ocup.%"]
    widths = [3.1, 1.35, 1.35, 1.35, 1.35, 1.25, 1.15, 1.15]
    body = [[r["operador"], r["logueo"], r["no_disponible"], r["pct_no_disponible"],
             r["atendidas"], r["cierres"], r["utilizacion"], r["ocupacion"]]
            for r in rows]

    # Keep the table inside the slide even with many operators
    n = len(body) + 1
    avail = 6.85 - 2.55
    row_h = min(0.28, max(0.16, (avail - 0.32) / max(len(body), 1)))
    fs = 9 if row_h >= 0.24 else (8 if row_h >= 0.19 else 7)

    _table(s, (13.333 - sum(widths)) / 2, 2.55, widths, headers, body,
           row_h=row_h, header_h=0.32, font_size=fs, header_size=9,
           bold_last=True)


# ======================================================================
# Entry point
# ======================================================================

def generate_plan_medico_pptx(
    period: str,
    headline: dict[str, str],
    section1: dict[str, Any],
    section2: dict[str, Any],
    section3a: dict[str, Any],
    closures: dict[str, Any] | None = None,
    agents: dict[str, Any] | None = None,
) -> bytes:
    """Build the Plan Medico deck and return it as bytes."""
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 4 + (1 if closures else 0) + (1 if agents else 0)

    _slide_cover(prs, period, headline)

    _slide_section(
        prs, period, 2, total, "Seccion 1 - Plan Medico Total",
        section1.get("subtitle", ""), section1["kpis"], section1.get("chart"))

    _slide_section(
        prs, period, 3, total, "Seccion 2 - Turnos Plan Medico",
        section2.get("subtitle", ""), section2["kpis"], section2.get("chart"),
        table=section2.get("table"))

    _slide_section(
        prs, period, 4, total,
        "Seccion 3a - Plan Medico Consultas: Indicadores operativos",
        section3a.get("subtitle", ""), section3a["kpis"], section3a.get("chart"))

    page = 5
    if closures:
        _slide_closures(prs, period, page, total,
                        closures.get("subtitle", ""), closures.get("chart"))
        page += 1

    if agents:
        _slide_agents(prs, period, page, total,
                      agents["summary"], agents["rows"])

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()
