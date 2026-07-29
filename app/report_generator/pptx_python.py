"""PPTX report generator using python-pptx (pure Python).

Generates the Contact Center monthly report matching the HA visual style.
No Node.js dependency — suitable for Streamlit Cloud deployment.
"""

from __future__ import annotations

import io
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

import pandas as pd


# ======================================================================
# Colors
# ======================================================================
DARK_NAVY   = RGBColor(0x1B, 0x3A, 0x5C)
MEDIUM_BLUE = RGBColor(0x5B, 0x9B, 0xD5)
GREEN       = RGBColor(0x4C, 0xAF, 0x50)
RED         = RGBColor(0xE7, 0x4C, 0x3C)
LIGHT_GRAY  = RGBColor(0xF5, 0xF6, 0xF8)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
TEXT_DARK   = RGBColor(0x2C, 0x3E, 0x50)
TEXT_GRAY   = RGBColor(0x7F, 0x8C, 0x8D)
BORDER_GRAY = RGBColor(0xE0, 0xE0, 0xE0)

FONT_NAME = "Calibri"

# Slide dimensions (LAYOUT_WIDE = 13.33" x 7.5")
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


# ======================================================================
# Helper functions
# ======================================================================

def _add_textbox(slide, left, top, width, height, text,
                 font_size=10, bold=False, italic=False,
                 color=TEXT_DARK, alignment=PP_ALIGN.LEFT,
                 font_name=FONT_NAME):
    """Add a styled text box to a slide."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(text)
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def _add_rect(slide, left, top, width, height, fill_color=None, line_color=None):
    """Add a rectangle shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.shadow.inherit = False
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.fill.solid()
        shape.line.fill.fore_color.rgb = line_color
        shape.line.width = Pt(0.5)
    else:
        shape.line.fill.background()
    return shape


def _add_header_bar(slide, title_text, subtitle_text, period):
    """Add the dark navy header bar with green accent."""
    # Navy bar
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.65), fill_color=DARK_NAVY)
    # Green accent
    _add_rect(slide, 0, Inches(0.65), SLIDE_W, Inches(0.04), fill_color=GREEN)
    # Title
    _add_textbox(slide, Inches(0.5), Inches(0.12), Inches(7), Inches(0.4),
                 title_text, font_size=18, bold=True, color=WHITE)
    # Subtitle
    if subtitle_text:
        _add_textbox(slide, Inches(8), Inches(0.12), Inches(5), Inches(0.4),
                     subtitle_text, font_size=10, color=RGBColor(0xB0, 0xBE, 0xC5),
                     alignment=PP_ALIGN.RIGHT)


def _add_footer(slide, page_num, total_pages):
    """Add footer line with source and page number."""
    _add_rect(slide, Inches(0.5), Inches(6.9), Inches(12.3), Pt(1),
              fill_color=RGBColor(0xCC, 0xCC, 0xCC))
    _add_textbox(slide, Inches(0.5), Inches(6.95), Inches(6), Inches(0.3),
                 "Hospital Alemán · Fuente: Tecnovoz",
                 font_size=8, color=TEXT_GRAY)
    _add_textbox(slide, Inches(9), Inches(6.95), Inches(3.8), Inches(0.3),
                 f"Página {page_num} de {total_pages}",
                 font_size=8, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)


def _add_section_title(slide, text, y_inches):
    """Add section title with left accent bar."""
    _add_rect(slide, Inches(0.5), Inches(y_inches), Inches(0.06), Inches(0.35),
              fill_color=DARK_NAVY)
    _add_textbox(slide, Inches(0.7), Inches(y_inches), Inches(10), Inches(0.35),
                 text, font_size=14, bold=True, color=TEXT_DARK)


def _add_kpi_card(slide, x, y, w, label, value, variation=None, accent_color=None):
    """Add a KPI card with label, value and optional variation."""
    h = Inches(1.0)
    # Card border
    _add_rect(slide, Inches(x), Inches(y), Inches(w), h,
              fill_color=WHITE, line_color=BORDER_GRAY)
    # Accent line
    if accent_color:
        _add_rect(slide, Inches(x), Inches(y), Inches(w), Pt(3),
                  fill_color=accent_color)
    # Label
    _add_textbox(slide, Inches(x + 0.15), Inches(y + 0.1), Inches(w - 0.3), Inches(0.2),
                 label.upper(), font_size=8, color=TEXT_GRAY)
    # Value
    _add_textbox(slide, Inches(x + 0.15), Inches(y + 0.3), Inches(w - 0.3), Inches(0.4),
                 str(value), font_size=24, bold=True, color=DARK_NAVY,
                 alignment=PP_ALIGN.CENTER)
    # Variation
    if variation and variation != "—":
        var_color = GREEN if "▲" in str(variation) else RED
        _add_textbox(slide, Inches(x + 0.15), Inches(y + 0.72), Inches(w - 0.3), Inches(0.22),
                     str(variation), font_size=9, color=var_color,
                     alignment=PP_ALIGN.CENTER)


def _add_time_card(slide, x, y, w, label, value):
    """Add a time KPI card."""
    h = Inches(0.85)
    _add_rect(slide, Inches(x), Inches(y), Inches(w), h,
              fill_color=WHITE, line_color=BORDER_GRAY)
    _add_textbox(slide, Inches(x + 0.15), Inches(y + 0.08), Inches(w - 0.3), Inches(0.2),
                 label.upper(), font_size=8, color=TEXT_GRAY)
    _add_textbox(slide, Inches(x + 0.15), Inches(y + 0.3), Inches(w - 0.3), Inches(0.35),
                 str(value), font_size=22, bold=True, color=DARK_NAVY,
                 alignment=PP_ALIGN.CENTER)


def _add_chart_image(slide, image_bytes_or_path, x, y, w, h):
    """Add a chart image to the slide."""
    if isinstance(image_bytes_or_path, (str, Path)):
        p = Path(image_bytes_or_path)
        if p.exists():
            slide.shapes.add_picture(str(p), Inches(x), Inches(y), Inches(w), Inches(h))
    elif isinstance(image_bytes_or_path, bytes):
        stream = io.BytesIO(image_bytes_or_path)
        slide.shapes.add_picture(stream, Inches(x), Inches(y), Inches(w), Inches(h))


# ======================================================================
# Slide builders
# ======================================================================

def _build_cover(prs, period):
    """Slide 1: Cover page with Hospital Alemán + JCI logo."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Top bands
    _add_rect(slide, 0, 0, SLIDE_W, Inches(0.35), fill_color=DARK_NAVY)
    _add_rect(slide, 0, Inches(0.35), SLIDE_W, Inches(0.04), fill_color=GREEN)

    # Logo (apaisado — horizontal HA + JCI). Centered horizontally.
    # The logo is embedded as base64 in logo_data.py so it can never be lost
    # during GitHub uploads. Falls back to a file path, then to text.
    logo_stream = None
    try:
        from app.report_generator.logo_data import LOGO_HA_APAISADO_B64
        import base64
        logo_stream = io.BytesIO(base64.b64decode(LOGO_HA_APAISADO_B64))
    except Exception:
        logo_path = Path(__file__).parent.parent / "plugins" / "contact_center" / "templates" / "assets" / "logo_ha_apaisado.png"
        if logo_path.exists():
            logo_stream = str(logo_path)

    if logo_stream is not None:
        # Logo ratio ~4.95 (1978x400). Target width 7", height ~1.42".
        logo_w = 7.0
        logo_h = logo_w / 4.95
        logo_x = (13.333 - logo_w) / 2
        slide.shapes.add_picture(logo_stream, Inches(logo_x), Inches(1.3),
                                 Inches(logo_w), Inches(logo_h))
    else:
        # Fallback: text if logo missing
        _add_textbox(slide, Inches(2), Inches(1.8), Inches(9), Inches(1.2),
                     "Hospital Alemán", font_size=44, bold=True, color=DARK_NAVY,
                     alignment=PP_ALIGN.CENTER)

    # Title with accent
    _add_rect(slide, Inches(0.8), Inches(3.6), Inches(0.06), Inches(0.5),
              fill_color=GREEN)
    _add_textbox(slide, Inches(1.0), Inches(3.5), Inches(10), Inches(0.7),
                 "Productividad del Contact Center",
                 font_size=28, bold=True, color=DARK_NAVY)
    _add_textbox(slide, Inches(1.0), Inches(4.2), Inches(10), Inches(0.4),
                 "Análisis de Campañas", font_size=16, color=TEXT_DARK)
    _add_textbox(slide, Inches(1.0), Inches(4.6), Inches(10), Inches(0.3),
                 "Volumen de llamadas, nivel de atención y tiempos operativos",
                 font_size=11, color=TEXT_GRAY)

    # Period box
    _add_rect(slide, Inches(1.0), Inches(5.3), Inches(3.5), Inches(0.7),
              fill_color=LIGHT_GRAY, line_color=BORDER_GRAY)
    _add_textbox(slide, Inches(1.15), Inches(5.35), Inches(3), Inches(0.2),
                 "PERÍODO ANALIZADO", font_size=8, color=TEXT_GRAY)
    _add_textbox(slide, Inches(1.15), Inches(5.55), Inches(3), Inches(0.35),
                 period, font_size=14, bold=True, color=DARK_NAVY)

    # Source
    _add_textbox(slide, Inches(9.5), Inches(5.35), Inches(3), Inches(0.2),
                 "FUENTE", font_size=8, color=TEXT_GRAY, alignment=PP_ALIGN.RIGHT)
    _add_textbox(slide, Inches(9.5), Inches(5.55), Inches(3), Inches(0.35),
                 "Tecnovoz", font_size=14, bold=True, color=DARK_NAVY,
                 alignment=PP_ALIGN.RIGHT)

    # Bottom bands
    _add_rect(slide, 0, Inches(7.1), SLIDE_W, Inches(0.04), fill_color=GREEN)
    _add_rect(slide, 0, Inches(7.14), SLIDE_W, Inches(0.36), fill_color=DARK_NAVY)


def _add_big_kpi_card(slide, x, y, w, h, label, value, variation=None, accent_color=None):
    """Add a larger, more prominent KPI card for the main dashboard slide."""
    _add_rect(slide, Inches(x), Inches(y), Inches(w), Inches(h),
              fill_color=WHITE, line_color=BORDER_GRAY)
    if accent_color:
        _add_rect(slide, Inches(x), Inches(y), Inches(w), Pt(5),
                  fill_color=accent_color)
    # Label
    _add_textbox(slide, Inches(x + 0.2), Inches(y + 0.2), Inches(w - 0.4), Inches(0.35),
                 label.upper(), font_size=13, color=TEXT_GRAY,
                 alignment=PP_ALIGN.CENTER)
    # Big value
    _add_textbox(slide, Inches(x + 0.2), Inches(y + 0.6), Inches(w - 0.4), Inches(h - 1.2),
                 str(value), font_size=54, bold=True, color=DARK_NAVY,
                 alignment=PP_ALIGN.CENTER)
    # Variation
    if variation and variation != "—":
        var_color = GREEN if "▲" in str(variation) else RED
        _add_textbox(slide, Inches(x + 0.2), Inches(y + h - 0.55), Inches(w - 0.4), Inches(0.35),
                     str(variation), font_size=14, bold=True, color=var_color,
                     alignment=PP_ALIGN.CENTER)


def _build_general_data(prs, period, kpis, variations):
    """Slide 2: General data with 5 large KPI cards in 2 rows.

    Row 1 (top): Recibidas | Atendidas
    Row 2 (bottom): Prom. Diario Recibidas | Prom. Diario Atendidas | Nivel de Atención
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = f"Productividad del Contact Center · {period}"
    _add_header_bar(slide, "Datos Generales", subtitle, period)
    _add_footer(slide, 2, 19)
    _add_section_title(slide, "Indicadores principales del mes", 0.9)

    _add_textbox(slide, Inches(0.7), Inches(1.35), Inches(12), Inches(0.3),
                 "Variaciones calculadas respecto al mes anterior.",
                 font_size=10, italic=True, color=TEXT_GRAY)

    # === ROW 1: Recibidas + Atendidas (large, centered) ===
    # Two big cards centered horizontally
    row1_h = 2.4
    row1_y = 1.85
    card_w = 5.5
    row1_gap = 0.4
    total_w = card_w * 2 + row1_gap
    row1_start = (13.333 - total_w) / 2  # Center in slide

    _add_big_kpi_card(slide, row1_start, row1_y, card_w, row1_h,
                      "Recibidas", kpis.get("recibidas", "—"),
                      variations.get("recibidas", ""), DARK_NAVY)
    _add_big_kpi_card(slide, row1_start + card_w + row1_gap, row1_y, card_w, row1_h,
                      "Atendidas", kpis.get("atendidas", "—"),
                      variations.get("atendidas", ""), MEDIUM_BLUE)

    # === ROW 2: Prom. Rec + Prom. Att + NA (3 medium cards, centered) ===
    row2_h = 2.2
    row2_y = 4.45
    card_w2 = 3.9
    row2_gap = 0.3
    total_w2 = card_w2 * 3 + row2_gap * 2
    row2_start = (13.333 - total_w2) / 2

    _add_big_kpi_card(slide, row2_start, row2_y, card_w2, row2_h,
                      "Prom. Diario Recibidas", kpis.get("promedio_recibidas", "—"),
                      None, TEXT_GRAY)
    _add_big_kpi_card(slide, row2_start + card_w2 + row2_gap, row2_y, card_w2, row2_h,
                      "Prom. Diario Atendidas", kpis.get("promedio_atendidas", "—"),
                      None, TEXT_GRAY)
    _add_big_kpi_card(slide, row2_start + (card_w2 + row2_gap) * 2, row2_y, card_w2, row2_h,
                      "Nivel de Atención", kpis.get("nivel_atencion", "—"),
                      variations.get("nivel_atencion", ""), GREEN)


def _build_campaign_slide(prs, name, kpis, variations, chart_path,
                           page_num, period, is_all=False):
    """Build a campaign slide with KPIs + time cards + chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = f"Productividad del Contact Center · {period}"
    title = name if is_all else f"Campaña: {name}"
    _add_header_bar(slide, title, subtitle, period)
    _add_footer(slide, page_num, 19)
    _add_section_title(slide, "Indicadores y distribución diaria", 0.9)

    cardW = 2.3
    gap = 0.15
    startX = 0.5

    cards = [
        ("Recibidas", kpis.get("recibidas", "—"), variations.get("recibidas", ""), DARK_NAVY),
        ("Atendidas", kpis.get("atendidas", "—"), variations.get("atendidas", ""), MEDIUM_BLUE),
        ("Prom. Recibidas", kpis.get("promedio_recibidas", "—"), None, None),
        ("Prom. Atendidas", kpis.get("promedio_atendidas", "—"), None, None),
        ("Nivel de Atención", kpis.get("nivel_atencion", "—"), variations.get("nivel_atencion", ""), GREEN),
    ]
    for i, (label, value, var, accent) in enumerate(cards):
        _add_kpi_card(slide, startX + i * (cardW + gap), 1.4, cardW,
                      label, value, var, accent)

    # Time cards
    timeW = 3.8
    _add_time_card(slide, startX, 2.6, timeW,
                   "Conversación", kpis.get("tiempo_conversacion", "—"))
    _add_time_card(slide, startX + timeW + gap, 2.6, timeW,
                   "Demora", kpis.get("tiempo_demora", "—"))
    _add_time_card(slide, startX + (timeW + gap) * 2, 2.6, timeW,
                   "Abandono", kpis.get("tiempo_abandono", "—"))

    # Chart
    if chart_path:
        _add_chart_image(slide, chart_path, 0.3, 3.7, 12.5, 3.0)


def _build_chart_slide(prs, title, section, chart_path, page_num, period):
    """Build a slide with just a header and a chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = f"Productividad del Contact Center · {period}"
    _add_header_bar(slide, title, subtitle, period)
    _add_footer(slide, page_num, 19)
    _add_section_title(slide, section, 0.9)
    if chart_path:
        _add_chart_image(slide, chart_path, 0.3, 1.4, 12.5, 5.2)


def _build_dual_chart_slide(prs, title, section, chart_left, chart_right,
                              page_num, period, footnote=None):
    """Build a slide with two charts side by side (e.g., campaign analysis)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = f"Productividad del Contact Center · {period}"
    _add_header_bar(slide, title, subtitle, period)
    _add_footer(slide, page_num, 19)
    _add_section_title(slide, section, 0.9)
    if chart_left:
        _add_chart_image(slide, chart_left, 0.2, 1.4, 6.5, 4.8)
    if chart_right:
        _add_chart_image(slide, chart_right, 6.8, 1.4, 6.0, 4.8)
    if footnote:
        # Light grey note box at the bottom, like the original report
        _add_rect(slide, Inches(0.5), Inches(6.25), Inches(12.3), Inches(0.5),
                  fill_color=LIGHT_GRAY)
        _add_textbox(slide, Inches(0.7), Inches(6.33), Inches(12.0), Inches(0.4),
                     footnote, font_size=9, italic=True, color=TEXT_GRAY)


def _build_skill_table_slide(prs, skill_table, page_num, period):
    """Build the skill detail table slide."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = f"Productividad del Contact Center · {period}"
    _add_header_bar(slide, "Análisis de Habilidades", subtitle, period)
    _add_footer(slide, page_num, 19)
    _add_section_title(slide, "Detalle por habilidad — volumen, atención y tiempos promedio", 0.9)

    if not skill_table:
        return

    headers = ["Habilidad", "Recibidas", "Atendidas", "NA",
               "Conversación", "Demora", "Abandono"]
    n_rows = len(skill_table) + 1
    n_cols = len(headers)

    from pptx.util import Inches, Pt

    table_shape = slide.shapes.add_table(
        n_rows, n_cols,
        Inches(0.5), Inches(1.4), Inches(12.3), Inches(min(n_rows * 0.32, 5.0))
    )
    table = table_shape.table

    # Column widths
    col_widths = [Inches(2.5), Inches(1.4), Inches(1.4), Inches(1.2),
                  Inches(1.5), Inches(1.5), Inches(1.5)]
    for i, w in enumerate(col_widths):
        table.columns[i].width = w

    # Header row
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(8)
            paragraph.font.bold = True
            paragraph.font.color.rgb = WHITE
            paragraph.font.name = FONT_NAME
            paragraph.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY

    # Data rows
    for i, sk in enumerate(skill_table):
        row_data = [sk["name"], sk["recibidas"], sk["atendidas"],
                    sk["na"], sk["conversacion"], sk["demora"], sk["abandono"]]
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE

        for j, val in enumerate(row_data):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(7)
                paragraph.font.name = FONT_NAME
                paragraph.alignment = PP_ALIGN.RIGHT if j > 0 else PP_ALIGN.LEFT

                # Color-code NA column
                if j == 3:
                    try:
                        na_val = float(str(val).replace(",", ".").replace("%", ""))
                        if na_val < 85:
                            paragraph.font.color.rgb = RED
                        elif na_val >= 95:
                            paragraph.font.color.rgb = GREEN
                        else:
                            paragraph.font.color.rgb = TEXT_DARK
                    except ValueError:
                        paragraph.font.color.rgb = TEXT_DARK


def _fill_annex_table(table, rows_data, headers):
    """Fill a table with header + data rows + optional total row.

    rows_data: list of [fecha, recibidas, atendidas, na] lists.
               The last row may be a total (rendered bold).
    """
    # Header
    for j, header in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = header
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(8)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY

    # Data rows
    for i, row in enumerate(rows_data):
        is_total = row[0] == "Total general"
        bg = LIGHT_GRAY if (i % 2 == 0 or is_total) else WHITE
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(8)
                p.font.bold = bool(is_total)
                p.font.name = FONT_NAME
                if is_total:
                    p.font.color.rgb = DARK_NAVY
                p.alignment = PP_ALIGN.CENTER if j == 0 else PP_ALIGN.RIGHT


def _build_annex_daily_table(prs, campaign_name, daily_rows, page_num, total_pages, period):
    """Build an annex slide with daily productivity table for one campaign.

    daily_rows: list of dicts with keys: fecha, recibidas, atendidas, na

    For months with many days (>16), the table is split into TWO
    side-by-side tables so everything fits inside the slide.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = f"Productividad del Contact Center · {period}"
    _add_header_bar(slide, f"Anexo — {campaign_name}", subtitle, period)
    _add_footer(slide, page_num, total_pages)
    _add_section_title(slide, f"Productividad diaria — {campaign_name}", 0.9)

    if not daily_rows:
        _add_textbox(slide, Inches(0.7), Inches(1.5), Inches(10), Inches(0.4),
                     "Sin datos disponibles", font_size=12, color=TEXT_GRAY, italic=True)
        return

    headers = ["Fecha", "Recib.", "Atend.", "NA"]

    # Compute totals
    total_rec = 0
    total_att = 0
    for row in daily_rows:
        try:
            total_rec += int(str(row["recibidas"]).replace(".", "").replace(",", ""))
            total_att += int(str(row["atendidas"]).replace(".", "").replace(",", ""))
        except ValueError:
            pass
    na_total = (total_att / total_rec * 100) if total_rec > 0 else 0
    total_row = ["Total general",
                 f"{total_rec:,}".replace(",", "."),
                 f"{total_att:,}".replace(",", "."),
                 f"{na_total:.2f}%".replace(".", ",")]

    # Convert dicts to lists
    data_rows = [[r["fecha"], r["recibidas"], r["atendidas"], r["na"]] for r in daily_rows]

    SPLIT_THRESHOLD = 16  # if more days than this, use two columns

    if len(data_rows) <= SPLIT_THRESHOLD:
        # Single table, centered
        all_rows = data_rows + [total_row]
        n_rows = len(all_rows) + 1  # +header
        table_w = 7.0
        table_x = (13.333 - table_w) / 2
        row_h = 0.28
        table_h = min(n_rows * row_h, 5.3)

        shape = slide.shapes.add_table(n_rows, 4, Inches(table_x), Inches(1.4),
                                       Inches(table_w), Inches(table_h))
        for i, w in enumerate([Inches(1.75)] * 4):
            shape.table.columns[i].width = w
        _fill_annex_table(shape.table, all_rows, headers)
    else:
        # Split into two side-by-side tables
        mid = (len(data_rows) + 1) // 2
        left_rows = data_rows[:mid]
        right_rows = data_rows[mid:] + [total_row]  # total goes at bottom of right table

        # Make both tables same number of rows for visual balance
        max_len = max(len(left_rows), len(right_rows))
        # Pad left table with empty rows if needed
        while len(left_rows) < max_len:
            left_rows.append(["", "", "", ""])

        table_w = 5.6
        gap = 0.5
        left_x = (13.333 - table_w * 2 - gap) / 2
        right_x = left_x + table_w + gap
        row_h = 0.26
        n_rows = max_len + 1  # +header
        table_h = min(n_rows * row_h, 5.3)

        col_w = [Inches(1.4)] * 4

        # Left table
        left_shape = slide.shapes.add_table(n_rows, 4, Inches(left_x), Inches(1.4),
                                            Inches(table_w), Inches(table_h))
        for i, w in enumerate(col_w):
            left_shape.table.columns[i].width = w
        _fill_annex_table(left_shape.table, left_rows, headers)

        # Right table
        right_shape = slide.shapes.add_table(n_rows, 4, Inches(right_x), Inches(1.4),
                                             Inches(table_w), Inches(table_h))
        for i, w in enumerate(col_w):
            right_shape.table.columns[i].width = w
        _fill_annex_table(right_shape.table, right_rows, headers)


# ======================================================================
# Main generator function
# ======================================================================

def _build_outbound_slide(prs, outbound, chart_images, page_num, total_pages, period):
    """Slide: Llamadas Salientes with KPIs + result distribution + daily chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = f"Productividad del Contact Center · {period}"
    _add_header_bar(slide, "Llamadas Salientes", subtitle, period)
    _add_footer(slide, page_num, total_pages)
    _add_section_title(slide, "Resumen de gestión de llamadas salientes", 0.9)

    # KPI cards. Rotaciones AM and Solo Operadores AM are left BLANK on purpose:
    # those figures come from the supervisors' cancellation registry, so the
    # user fills them in manually in PowerPoint.
    total = outbound.get("total", 0)

    def _fmt(n):
        return f"{int(n):,}".replace(",", ".")

    _add_kpi_card(slide, 0.5, 1.4, 3.8, "Total Llamadas Salientes", _fmt(total), None, DARK_NAVY)
    _add_kpi_card(slide, 4.6, 1.4, 3.8, "Rotaciones AM", "", None, MEDIUM_BLUE)
    _add_kpi_card(slide, 8.7, 1.4, 3.8, "Solo Operadores AM", "", None, GREEN)

    # Hint that these two are filled manually
    _add_textbox(slide, 4.6 * 914400, int(2.42 * 914400), int(7.9 * 914400), int(0.25 * 914400),
                 "Completar manualmente — fuente: registro de cancelaciones de supervisores",
                 font_size=8, italic=True, color=TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    # Two charts side by side
    if chart_images.get("outbound_result"):
        _add_chart_image(slide, chart_images["outbound_result"], 0.3, 2.7, 6.2, 4.0)
    if chart_images.get("outbound_daily"):
        _add_chart_image(slide, chart_images["outbound_daily"], 6.7, 2.7, 6.3, 4.0)


def _build_monthly_trend_slide(prs, trend_records, chart_path, page_num, total_pages, period):
    """Slide: Evolución Mensual with a table + trend chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = f"Productividad del Contact Center · {period}"
    _add_header_bar(slide, "Evolución Mensual", subtitle, period)
    _add_footer(slide, page_num, total_pages)
    _add_section_title(slide, "Tendencia mensual del año", 0.9)

    if not trend_records:
        _add_textbox(slide, Inches(0.7), Inches(1.5), Inches(10), Inches(0.4),
                     "Sin datos históricos disponibles", font_size=12,
                     color=TEXT_GRAY, italic=True)
        return

    # Compact table at top
    n_months = len(trend_records)
    headers = ["", *[r["month_name"] for r in trend_records]]
    rows_spec = [
        ("Recibidas", [f"{r['recibidas']:,}".replace(",", ".") for r in trend_records]),
        ("Atendidas", [f"{r['atendidas']:,}".replace(",", ".") for r in trend_records]),
        ("Nivel de Atención", [f"{r['nivel_atencion']:.1f}%".replace(".", ",") for r in trend_records]),
    ]

    n_cols = n_months + 1
    table_w = min(1.2 + n_months * 2.0, 12.3)
    table_x = (13.333 - table_w) / 2
    shape = slide.shapes.add_table(4, n_cols, Inches(table_x), Inches(1.35),
                                   Inches(table_w), Inches(1.4))
    table = shape.table

    # Header
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.name = FONT_NAME
            p.font.color.rgb = WHITE
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_NAVY

    # Rows
    for i, (label, values) in enumerate(rows_spec):
        cell = table.cell(i + 1, 0)
        cell.text = label
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(10)
            p.font.bold = True
            p.font.name = FONT_NAME
            p.alignment = PP_ALIGN.LEFT
        cell.fill.solid()
        cell.fill.fore_color.rgb = LIGHT_GRAY
        for j, val in enumerate(values):
            c = table.cell(i + 1, j + 1)
            c.text = val
            for p in c.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = FONT_NAME
                p.alignment = PP_ALIGN.CENTER
            c.fill.solid()
            c.fill.fore_color.rgb = WHITE

    # Chart below the table
    if chart_path:
        _add_chart_image(slide, chart_path, 0.5, 3.0, 12.3, 3.7)


# ======================================================================
# Main generator function
# ======================================================================

def _build_skills_reference_annex(prs, skills_reference, page_num, total_pages, period):
    """Annex slide: reference table mapping each skill to its campaign.

    skills_reference: list of {"skill": ..., "campaign": ...}
    Rendered in up to 3 side-by-side columns so all 19+ skills fit.
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    subtitle = f"Productividad del Contact Center · {period}"
    _add_header_bar(slide, "Anexo — Referencia de Habilidades", subtitle, period)
    _add_footer(slide, page_num, total_pages)
    _add_section_title(slide, "Habilidades por campaña", 0.9)

    if not skills_reference:
        _add_textbox(slide, Inches(0.7), Inches(1.5), Inches(10), Inches(0.4),
                     "Sin datos disponibles", font_size=12, color=TEXT_GRAY, italic=True)
        return

    rows = [[s["skill"], s["campaign"]] for s in skills_reference]
    headers = ["Habilidad", "Campaña"]

    n = len(rows)
    n_cols_layout = 1 if n <= 14 else (2 if n <= 28 else 3)
    per_col = -(-n // n_cols_layout)
    chunks = [rows[i * per_col:(i + 1) * per_col] for i in range(n_cols_layout)]
    chunks = [c for c in chunks if c]

    table_w = 3.9 if n_cols_layout == 3 else (5.6 if n_cols_layout == 2 else 6.5)
    gap = 0.35
    total_w = table_w * len(chunks) + gap * (len(chunks) - 1)
    start_x = (13.333 - total_w) / 2
    max_len = max(len(c) for c in chunks)
    row_h = 0.26
    table_h = min((max_len + 1) * row_h, 5.3)

    for idx, chunk in enumerate(chunks):
        padded = list(chunk) + [["", ""]] * (max_len - len(chunk))
        x = start_x + idx * (table_w + gap)
        shape = slide.shapes.add_table(max_len + 1, 2, Inches(x), Inches(1.4),
                                       Inches(table_w), Inches(table_h))
        table = shape.table
        table.columns[0].width = Inches(table_w * 0.55)
        table.columns[1].width = Inches(table_w * 0.45)

        for j, h in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = h
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(9)
                p.font.bold = True
                p.font.color.rgb = WHITE
                p.font.name = FONT_NAME
                p.alignment = PP_ALIGN.CENTER
            cell.fill.solid()
            cell.fill.fore_color.rgb = DARK_NAVY

        for i, row in enumerate(padded):
            bg = LIGHT_GRAY if i % 2 == 0 else WHITE
            for j, val in enumerate(row):
                cell = table.cell(i + 1, j)
                cell.text = str(val)
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
                for p in cell.text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.name = FONT_NAME
                    p.alignment = PP_ALIGN.LEFT


def generate_pptx_report(
    period: str,
    global_kpis: dict[str, str],
    global_variations: dict[str, str],
    campaign_data: list[dict[str, Any]],
    skill_table: list[dict[str, str]],
    chart_images: dict[str, Any],
    annexes: list[dict[str, Any]] | None = None,
    outbound: dict[str, Any] | None = None,
    monthly_trend: list[dict[str, Any]] | None = None,
    donut_footnote: str | None = None,
    skills_reference: list[dict[str, str]] | None = None,
) -> bytes:
    """Generate a complete PPTX report and return it as bytes.

    Parameters
    ----------
    period : str
        e.g. "Mayo 2026"
    global_kpis : dict
        Formatted KPI values keyed by KPI id.
    global_variations : dict
        Formatted variation strings.
    campaign_data : list of dict
        Each dict has: name, kpis, variations, chart_path
    skill_table : list of dict
        Each dict has: name, recibidas, atendidas, na, conversacion, demora, abandono
    chart_images : dict
        Chart id → file path or bytes
    annexes : list of dict, optional
        Daily productivity tables per campaign. Each dict has:
          - campaign_name: str
          - daily_rows: list of {fecha, recibidas, atendidas, na}

    Returns
    -------
    bytes
        The PPTX file content.
    """
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # Calculate total pages upfront
    annexes = annexes or []
    total_pages = 4 + len(campaign_data) + 1 + len(annexes)  # cover+datos+weekday+campanalysis + N campaigns + skill_table + annexes

    # 1. Cover
    _build_cover(prs, period)

    # 2. General Data
    _build_general_data(prs, period, global_kpis, global_variations)

    # 3. Weekday distribution
    _build_chart_slide(prs, "Distribución por Día de Semana",
                       "Comportamiento por día — todas las campañas",
                       chart_images.get("weekday_distribution"),
                       3, period)

    # 4. Campaign analysis (dual charts)
    _build_dual_chart_slide(prs, "Análisis de Campañas",
                            "Volumen y participación por campaña",
                            chart_images.get("campaign_volume"),
                            chart_images.get("campaign_share"),
                            4, period,
                            footnote=donut_footnote)

    # 5+. Individual campaigns
    for i, camp in enumerate(campaign_data):
        _build_campaign_slide(
            prs, camp["name"], camp["kpis"], camp.get("variations", {}),
            camp.get("chart_path"), 5 + i, period
        )

    page = 5 + len(campaign_data)

    # Monthly trend (evolución mensual)
    if monthly_trend:
        _build_monthly_trend_slide(prs, monthly_trend,
                                   chart_images.get("monthly_evolution"),
                                   page, total_pages, period)
        page += 1

    # Outbound calls (llamadas salientes)
    if outbound:
        _build_outbound_slide(prs, outbound, chart_images, page, total_pages, period)
        page += 1

    # Top 10 skills chart
    if chart_images.get("skill_volume_top10"):
        _build_chart_slide(prs, "Análisis de Habilidades — Top 10",
                           "Top 10 habilidades por volumen de llamadas",
                           chart_images["skill_volume_top10"], page, period)
        page += 1

    # Skill detail table
    _build_skill_table_slide(prs, skill_table, page, period)
    skill_page = page

    # Annexes: one slide per campaign with daily table
    for i, annex in enumerate(annexes):
        _build_annex_daily_table(
            prs,
            annex["campaign_name"],
            annex["daily_rows"],
            skill_page + 1 + i,
            total_pages,
            period,
        )

    # Final annex: skill → campaign reference table
    if skills_reference:
        _build_skills_reference_annex(prs, skills_reference,
                                      skill_page + 1 + len(annexes),
                                      total_pages, period)

    # Save to bytes
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer.getvalue()
